import customtkinter as ctk
from tkinter import messagebox
from huggingface_hub import InferenceClient
from pptx import Presentation
from PIL import Image
import datetime
import threading

HF_TOKEN = "Add_huggingface_token"
MODEL_NAME = "deepseek-ai/DeepSeek-V3-0324"

client = InferenceClient(api_key=HF_TOKEN, provider="auto")

ctk.set_appearance_mode("dark")

ORANGE = "#ff9900"
GOLD = "#ffd166"
BG = "#070b14"
CARD = "#101722"
CHAT_BG = "#090d16"
BROWN = "#2b1b12"
WHITE = "#ffffff"

mode = "general"
spinner_running = False
spinner_index = 0
spinner_chars = ["◐", "◓", "◑", "◒"]

chat_file = "studentmate_chat_" + datetime.datetime.now().strftime("%Y%m%d_%H%M%S") + ".txt"

SYSTEM_PROMPT = """
You are StudentMate AI, a friendly assistant for students.

For college students, help with hackathons, project ideas, PPT content, tech stack, demo plan, pitch line, and idea scoring.

For school students, help with small projects, summaries, simple explanations, and tough questions in easy words.

Do not use markdown symbols like ###, **, or code blocks.
Talk naturally like a friendly student mentor.
Keep answers short, useful, and clear.
"""

messages = [{"role": "system", "content": SYSTEM_PROMPT}]


def save_chat(text):
    with open(chat_file, "a", encoding="utf-8") as file:
        file.write(text + "\n")


def clean_text(text):
    for bad in ["###", "**", "```"]:
        text = text.replace(bad, "")
    return text.strip()


def add_chat(sender, text):
    chat_box.configure(state="normal")

    if sender == "user":
        chat_box.insert("end", "\nYou\n", "user_title")
        chat_box.insert("end", text + "\n", "user_msg")

    elif sender == "bot":
        chat_box.insert("end", "\nStudentMate AI\n", "bot_title")
        chat_box.insert("end", text + "\n", "bot_msg")

    else:
        chat_box.insert("end", "\n" + text + "\n", "info")

    chat_box.configure(state="disabled")
    chat_box.see("end")


def start_spinner():
    global spinner_running
    spinner_running = True
    update_spinner()


def update_spinner():
    global spinner_index
    if spinner_running:
        status_label.configure(text=f"StudentMate AI thinking {spinner_chars[spinner_index]}")
        spinner_index = (spinner_index + 1) % len(spinner_chars)
        app.after(300, update_spinner)


def stop_spinner():
    global spinner_running
    spinner_running = False
    status_label.configure(text="")


def set_college():
    global mode
    mode = "college"
    show_options("college")


def set_school():
    global mode
    mode = "school"
    show_options("school")


def show_options(selected):
    for widget in options_frame.winfo_children():
        widget.destroy()

    if selected == "college":
        options = [
            ("🚀 Hackathon", "Give me national level hackathon project ideas"),
            ("💻 Project", "Help me build a college project"),
            ("📊 PPT", "Create PPT content for my project"),
        ]
    else:
        options = [
            ("📚 Project", "Suggest a simple school project"),
            ("📝 Summary", "Summarize this topic in simple points"),
            ("❓ Explain", "Explain this question in simple words"),
        ]

    for text, prompt in options:
        btn = ctk.CTkButton(
            options_frame,
            text=text,
            width=190,
            height=45,
            corner_radius=18,
            fg_color="#111111",
            hover_color="#2a2a2a",
            border_width=2,
            border_color=ORANGE,
            text_color=GOLD,
            font=("Arial", 14, "bold"),
            command=lambda p=prompt: insert_prompt(p)
        )
        btn.pack(side="left", padx=8, pady=5)


def insert_prompt(prompt):
    entry.delete(0, "end")
    entry.insert(0, prompt)


def send_message(event=None):
    user_text = entry.get().strip()

    if user_text == "":
        return

    entry.delete(0, "end")
    add_chat("user", user_text)
    save_chat("You: " + user_text)

    if mode == "college":
        final_text = "College student request: " + user_text
    elif mode == "school":
        final_text = "School student request: " + user_text
    else:
        final_text = user_text

    messages.append({"role": "user", "content": final_text})
    start_spinner()

    def run_ai():
        try:
            response = client.chat.completions.create(
                model=MODEL_NAME,
                messages=messages,
                max_tokens=350
            )

            reply = clean_text(response.choices[0].message.content)
            messages.append({"role": "assistant", "content": reply})

            app.after(0, stop_spinner)
            app.after(0, lambda: add_chat("bot", reply))
            save_chat("StudentMate AI: " + reply)

        except Exception as e:
            app.after(0, stop_spinner)
            app.after(0, lambda: add_chat("bot", "Error: " + str(e)))

    threading.Thread(target=run_ai, daemon=True).start()


def clear_chat():
    chat_box.configure(state="normal")
    chat_box.delete("1.0", "end")
    chat_box.configure(state="disabled")
    show_welcome()


def show_welcome():
    add_chat(
        "info",
        "✨ Welcome to StudentMate AI ✨\n"
        "Your smart assistant for projects, hackathons, summaries and learning.\n"
        "Choose College Student or School Student mode to begin."
    )


def make_ppt():
    topic = entry.get().strip()

    if topic == "":
        messagebox.showwarning("Input Needed", "Type your project idea or topic first.")
        return

    start_spinner()

    ppt_prompt = f"""
Create PPT content for this topic: {topic}

Use simple student-friendly language.
Do not use markdown symbols.

Create slide-wise content:
Slide 1: Title
Slide 2: Problem or Introduction
Slide 3: Objective
Slide 4: Proposed Solution or Main Content
Slide 5: Key Features or Main Points
Slide 6: Tech Stack or Materials Required
Slide 7: Workflow or Methodology
Slide 8: Benefits or Impact
Slide 9: Future Scope
Slide 10: Conclusion
"""

    def run_ppt():
        try:
            response = client.chat.completions.create(
                model=MODEL_NAME,
                messages=[
                    {"role": "system", "content": SYSTEM_PROMPT},
                    {"role": "user", "content": ppt_prompt}
                ],
                max_tokens=900
            )

            ppt_text = clean_text(response.choices[0].message.content)

            prs = Presentation()
            parts = ppt_text.split("Slide ")

            for part in parts:
                if part.strip() == "":
                    continue

                lines = part.strip().split("\n")
                title = "Slide " + lines[0].strip()
                body = "\n".join(lines[1:]).strip()

                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = title
                slide.placeholders[1].text = body

            ppt_name = "StudentMate_PPT_" + datetime.datetime.now().strftime("%Y%m%d_%H%M%S") + ".pptx"
            prs.save(ppt_name)

            app.after(0, stop_spinner)
            app.after(0, lambda: add_chat("bot", "PPT created successfully: " + ppt_name))
            app.after(0, lambda: messagebox.showinfo("PPT Created", "Saved as " + ppt_name))

        except Exception as e:
            app.after(0, stop_spinner)
            app.after(0, lambda: messagebox.showerror("Error", str(e)))

    threading.Thread(target=run_ppt, daemon=True).start()


app = ctk.CTk()
app.title("StudentMate AI")
app.geometry("1200x820")
app.minsize(1050, 740)
app.configure(fg_color=BG)

app.grid_rowconfigure(0, weight=1)
app.grid_columnconfigure(0, weight=1)

main = ctk.CTkFrame(
    app,
    fg_color=CARD,
    corner_radius=28,
    border_width=2,
    border_color=ORANGE
)
main.grid(row=0, column=0, padx=25, pady=25, sticky="nsew")

main.grid_rowconfigure(3, weight=1)
main.grid_columnconfigure(0, weight=1)

header = ctk.CTkFrame(main, fg_color="transparent")
header.grid(row=0, column=0, sticky="ew", padx=30, pady=(15, 5))
header.grid_columnconfigure(0, weight=1)

left_header = ctk.CTkFrame(header, fg_color="transparent")
left_header.grid(row=0, column=0, sticky="ew")

title = ctk.CTkLabel(
    left_header,
    text="StudentMate AI",
    text_color=ORANGE,
    font=("Arial", 46, "bold")
)
title.pack()

subtitle = ctk.CTkLabel(
    left_header,
    text="Your smart AI companion for projects, hackathons, summaries and learning",
    text_color=WHITE,
    font=("Arial", 16)
)
subtitle.pack(pady=4)

made = ctk.CTkLabel(
    left_header,
    text="Made by Tharun N",
    text_color=GOLD,
    font=("Arial", 14, "bold")
)
made.pack()

try:
    bot_img = ctk.CTkImage(
        dark_image=Image.open("bot.png"),
        size=(130, 130)
    )
    bot_label = ctk.CTkLabel(header, image=bot_img, text="")
    bot_label.grid(row=0, column=1, padx=25)
except:
    bot_label = ctk.CTkLabel(header, text="🤖", font=("Arial", 70), text_color=ORANGE)
    bot_label.grid(row=0, column=1, padx=25)

mode_frame = ctk.CTkFrame(main, fg_color="transparent")
mode_frame.grid(row=1, column=0, pady=8)

college_btn = ctk.CTkButton(
    mode_frame,
    text="🎓 College Student",
    width=300,
    height=65,
    corner_radius=20,
    fg_color="#1a1208",
    hover_color="#3a2208",
    border_width=2,
    border_color=ORANGE,
    text_color=GOLD,
    font=("Arial", 17, "bold"),
    command=set_college
)
college_btn.pack(side="left", padx=15)

school_btn = ctk.CTkButton(
    mode_frame,
    text="📖 School Student",
    width=300,
    height=65,
    corner_radius=20,
    fg_color="#111827",
    hover_color="#1f2937",
    border_width=2,
    border_color="#5a5a5a",
    text_color=GOLD,
    font=("Arial", 17, "bold"),
    command=set_school
)
school_btn.pack(side="left", padx=15)

options_frame = ctk.CTkFrame(main, fg_color="transparent")
options_frame.grid(row=2, column=0, pady=4)

chat_frame = ctk.CTkFrame(
    main,
    fg_color=CHAT_BG,
    corner_radius=22,
    border_width=2,
    border_color=ORANGE
)
chat_frame.grid(row=3, column=0, padx=35, pady=10, sticky="nsew")
chat_frame.grid_rowconfigure(0, weight=1)
chat_frame.grid_columnconfigure(0, weight=1)

chat_box = ctk.CTkTextbox(
    chat_frame,
    fg_color=CHAT_BG,
    text_color=WHITE,
    font=("Arial", 14),
    corner_radius=18,
    wrap="word"
)
chat_box.grid(row=0, column=0, padx=10, pady=10, sticky="nsew")

chat_box.tag_config("user_title", foreground=GOLD)
chat_box.tag_config("user_msg", foreground=WHITE)
chat_box.tag_config("bot_title", foreground=ORANGE)
chat_box.tag_config("bot_msg", foreground=WHITE)
chat_box.tag_config("info", foreground=ORANGE)

chat_box.configure(state="disabled")

status_label = ctk.CTkLabel(
    main,
    text="",
    text_color=ORANGE,
    font=("Arial", 14, "bold")
)
status_label.grid(row=4, column=0, pady=(0, 3))

bottom = ctk.CTkFrame(main, fg_color="transparent")
bottom.grid(row=5, column=0, sticky="ew", padx=35, pady=(5, 20))
bottom.grid_columnconfigure(0, weight=1)

entry = ctk.CTkEntry(
    bottom,
    placeholder_text="Type your message here...",
    height=62,
    corner_radius=18,
    fg_color=BROWN,
    border_color="#777777",
    border_width=2,
    text_color=WHITE,
    font=("Arial", 18)
)
entry.grid(row=0, column=0, sticky="ew", padx=(0, 12))
entry.bind("<Return>", send_message)

send_btn = ctk.CTkButton(
    bottom,
    text="➤ Send",
    width=145,
    height=62,
    corner_radius=18,
    fg_color=GOLD,
    hover_color=ORANGE,
    text_color="black",
    font=("Arial", 17, "bold"),
    command=send_message
)
send_btn.grid(row=0, column=1, padx=5)

ppt_btn = ctk.CTkButton(
    bottom,
    text="📊 PPT",
    width=125,
    height=62,
    corner_radius=18,
    fg_color="#111111",
    hover_color="#2a2a2a",
    border_width=2,
    border_color=ORANGE,
    text_color=GOLD,
    font=("Arial", 16, "bold"),
    command=make_ppt
)
ppt_btn.grid(row=0, column=2, padx=5)

clear_btn = ctk.CTkButton(
    bottom,
    text="Clear",
    width=115,
    height=62,
    corner_radius=18,
    fg_color="#111111",
    hover_color="#2a2a2a",
    border_width=2,
    border_color="#555555",
    text_color=GOLD,
    font=("Arial", 16, "bold"),
    command=clear_chat
)
clear_btn.grid(row=0, column=3, padx=5)

show_welcome()
app.mainloop()
