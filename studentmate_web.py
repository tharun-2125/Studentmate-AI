import streamlit as st
from huggingface_hub import InferenceClient
from pptx import Presentation
import datetime, json, hashlib, os

HF_TOKEN = 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
MODEL = "deepseek-ai/DeepSeek-V3-0324"

client = InferenceClient(
    api_key=st.secrets["HF_TOKEN"],
    provider="auto"
)

st.set_page_config(page_title="StudentMate AI", page_icon="🤖", layout="wide")

# ---------- CSS DESIGN ----------
st.markdown("""
<style>
.stApp {
    background: radial-gradient(circle at top, #101722 0%, #070b14 55%, #02040a 100%);
    color: white;
}
.block-container {
    padding-top: 20px;
    max-width: 1200px;
}
.main-card, .section-card {
    background: rgba(10, 15, 25, 0.92);
    border: 1.8px solid #ff9900;
    border-radius: 22px;
    padding: 24px;
    box-shadow: 0 0 25px rgba(255,153,0,0.25);
    margin-bottom: 18px;
}
.hero-title {
    font-size: 70px;
    font-weight: 900;
    color: #ff9900;
    text-shadow: 0 0 18px rgba(255,153,0,0.9);
    margin-bottom: 5px;
}
.subtitle {
    font-size: 24px;
    color: white;
}
.made {
    color: #ffd166;
    font-weight: bold;
    font-size: 18px;
}
.stButton button {
    height: 58px;
    border-radius: 16px;
    background: linear-gradient(135deg, #ffd166, #ff9900);
    color: black;
    font-weight: 800;
    border: 1px solid #ff9900;
}
.stTextInput input {
    background-color: #090d16;
    color: white;
    border: 1.5px solid #ff9900;
    border-radius: 16px;
    font-size: 17px;
}

div.stButton > button:first-child {
    background: linear-gradient(135deg, #ffd166, #ff9900) !important;
    color: black !important;
    font-weight: 800 !important;
    border-radius: 16px !important;
    height: 58px !important;
    border: 1px solid #ff9900 !important;
}

button[kind="secondary"] {
    background: linear-gradient(135deg, #ffd166, #ff9900) !important;
    color: black !important;
}

button[kind="primary"] {
    background: linear-gradient(135deg, #ffd166, #ff9900) !important;
    color: black !important;
}


.chat-box {
    background: #090d16;
    border: 1.5px solid #2f3a4d;
    border-radius: 18px;
    min-height: 280px;
    padding: 25px;
    font-size: 18px;
}
.user-msg {
    background:#2b1b12;
    padding:12px 16px;
    border-radius:16px;
    margin:10px 0;
    width:fit-content;
    max-width:70%;
    color:white;
}
.bot-msg {
    background:#111827;
    padding:12px 16px;
    border-radius:16px;
    margin:10px 0 10px auto;
    width:fit-content;
    max-width:70%;
    color:white;
    border-left:4px solid #ff9900;
}
</style>
""", unsafe_allow_html=True)

# ---------- LOGIN ----------
def load_users():
    try:
        with open("users.json", "r") as f:
            return json.load(f)
    except:
        return {}

def save_users(users):
    with open("users.json", "w") as f:
        json.dump(users, f)

def hash_password(p):
    return hashlib.sha256(p.encode()).hexdigest()

users = load_users()

if "logged_in" not in st.session_state:
    st.session_state.logged_in = False

if "chat_history" not in st.session_state:
    st.session_state.chat_history = []

if "mode" not in st.session_state:
    st.session_state.mode = None

if "task" not in st.session_state:
    st.session_state.task = None

if not st.session_state.logged_in:
    st.markdown("<div class='main-card'>", unsafe_allow_html=True)
    st.markdown("<h1 style='text-align:center;color:#ff9900;'>🔐 StudentMate AI Login</h1>", unsafe_allow_html=True)
    st.markdown("<p style='text-align:center;'>Login to access your smart student assistant</p>", unsafe_allow_html=True)

    tab1, tab2 = st.tabs(["Login", "Register"])

    with tab1:
        username = st.text_input("Username")
        password = st.text_input("Password", type="password")

        if st.button("Login"):
            if username in users and users[username] == hash_password(password):
                st.session_state.logged_in = True
                st.session_state.username = username
                st.rerun()
            else:
                st.error("Invalid username or password")

    with tab2:
        new_user = st.text_input("Create Username")
        new_pass = st.text_input("Create Password", type="password")

        if st.button("Register"):
            if new_user == "" or new_pass == "":
                st.warning("Enter username and password")
            elif new_user in users:
                st.error("Username already exists")
            else:
                users[new_user] = hash_password(new_pass)
                save_users(users)
                st.success("Account created. Now login.")

    st.markdown("</div>", unsafe_allow_html=True)
    st.stop()

# ---------- AI ----------
SYSTEM_PROMPT = """
You are StudentMate AI, a friendly assistant for students.
Help college students with hackathons, projects, PPT content, tech stack, demo plan, pitch line and idea scoring.
Help school students with small projects, summaries, simple explanations and tough questions.
Do not use markdown symbols like ###, ** or code blocks.
Talk naturally like a friendly student mentor.
Keep answers short, useful and clear.
"""

def clean_text(text):
    for bad in ["###", "**", "```"]:
        text = text.replace(bad, "")
    return text.strip()

def ask_ai(prompt, max_tokens=450):
    response = client.chat.completions.create(
        model=MODEL,
        messages=[
            {"role": "system", "content": SYSTEM_PROMPT},
            {"role": "user", "content": prompt}
        ],
        max_tokens=max_tokens
    )
    return clean_text(response.choices[0].message.content)

def create_ppt(topic):
    prompt = f"""
Create PPT content for this topic: {topic}

Slide 1: Title
Slide 2: Introduction
Slide 3: Objective
Slide 4: Main Content or Solution
Slide 5: Key Features or Main Points
Slide 6: Tech Stack or Materials Required
Slide 7: Workflow or Methodology
Slide 8: Benefits or Impact
Slide 9: Future Scope
Slide 10: Conclusion
"""
    ppt_text = ask_ai(prompt, 900)

    prs = Presentation()
    parts = ppt_text.split("Slide ")

    for part in parts:
        if part.strip() == "":
            continue

        lines = part.strip().split("\n")
        title = "Slide " + lines[0].strip()
        body = "\n".join(lines[1:]).strip()

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = body

    ppt_name = "StudentMate_PPT_" + datetime.datetime.now().strftime("%Y%m%d_%H%M%S") + ".pptx"
    prs.save(ppt_name)
    return ppt_name

# ---------- HEADER ----------
top1, top2, top3 = st.columns([6, 1, 1])

with top1:
    st.markdown("<h2 style='color:#ff9900;'>🤖 StudentMate AI</h2>", unsafe_allow_html=True)

with top2:
    st.markdown(f"<p style='color:#ffd166;'>👨‍🎓 {st.session_state.username}</p>", unsafe_allow_html=True)

with top3:
    if st.button("Logout"):
        st.session_state.logged_in = False
        st.session_state.chat_history = []
        st.rerun()

st.markdown("<div class='main-card'>", unsafe_allow_html=True)

col1, col2 = st.columns([3, 1])

with col1:
    st.markdown("<div class='hero-title'>✨ StudentMate AI</div>", unsafe_allow_html=True)
    st.markdown("<div class='subtitle'>Your smart assistant for<br>projects, hackathons, summaries and learning.</div>", unsafe_allow_html=True)
    st.markdown("<br><div class='made'>— ✦ Made by Tharun N ✦ —</div>", unsafe_allow_html=True)

with col2:
    if os.path.exists("bot.png"):
        st.image("bot.png", width=230)
    else:
        st.markdown("<h1 style='font-size:120px;'>🤖</h1>", unsafe_allow_html=True)

st.markdown("</div>", unsafe_allow_html=True)

# ---------- MODE ----------
st.markdown("<div class='section-card'>", unsafe_allow_html=True)
st.markdown("<h3 style='text-align:center;'>👨‍🎓 Choose Your Mode</h3>", unsafe_allow_html=True)

m1, m2 = st.columns(2)

with m1:
    if st.button("🎓 College Student", use_container_width=True):
        st.session_state.mode = "College Student"
        st.session_state.task = None
        st.rerun()

with m2:
    if st.button("🏫 School Student", use_container_width=True):
        st.session_state.mode = "School Student"
        st.session_state.task = None
        st.rerun()

st.markdown("</div>", unsafe_allow_html=True)

# ---------- OPTIONS ONLY AFTER MODE ----------
# ---------- OPTIONS ONLY AFTER MODE ----------
if st.session_state.mode is not None:

    st.markdown("<div class='section-card'>", unsafe_allow_html=True)
    st.markdown("<h3 style='text-align:center;'>✨ Choose an Option</h3>", unsafe_allow_html=True)

    if st.session_state.mode == "College Student":
        options = ["🚀 Hackathon", "💻 Project", "📊 PPT", "🎤 Pitch", "⭐ Score"]
    else:
        options = ["📚 Project", "📝 Summary", "❓ Explain", "✏️ Short Answer", "🧪 Activity"]

    cols = st.columns(len(options))

    for i, opt in enumerate(options):

        with cols[i]:

            # Selected button style
            if st.session_state.task == opt:

                st.markdown(f"""
                <div style="
                background: linear-gradient(135deg,#ff9900,#ffd166);
                padding:14px;
                border-radius:16px;
                text-align:center;
                font-weight:bold;
                color:black;
                box-shadow:0 0 15px rgba(255,153,0,0.9);
                ">
                ✅ {opt}
                </div>
                """, unsafe_allow_html=True)

            else:
                if st.button(opt, use_container_width=True):
                    st.session_state.task = opt
                    st.rerun()

    st.markdown("</div>", unsafe_allow_html=True)
# ---------- CHAT ----------
st.markdown("<div class='section-card'>", unsafe_allow_html=True)
st.markdown("<h3 style='text-align:center;color:#ffd166;'>💬 Chat with StudentMate AI</h3>", unsafe_allow_html=True)

if len(st.session_state.chat_history) == 0:
    st.markdown("""
    <div class='chat-box'>
    <h3 style='color:#ff9900;text-align:center;'>✨ Welcome to StudentMate AI ✨</h3>
    <p style='text-align:center;'>Choose mode and option, then start chatting!</p>
    </div>
    """, unsafe_allow_html=True)
else:
    st.markdown("<div class='chat-box'>", unsafe_allow_html=True)

    for msg in st.session_state.chat_history:
        if msg["role"] == "user":
            st.markdown(f"<div class='user-msg'>You: {msg['text']}</div>", unsafe_allow_html=True)
        else:
            st.markdown(f"<div class='bot-msg'>StudentMate AI: {msg['text']}</div>", unsafe_allow_html=True)

    st.markdown("</div>", unsafe_allow_html=True)

with st.form("chat_form", clear_on_submit=True):
    user_input = st.text_input(
        "",
        placeholder="Type your question or project idea here..."
    )

    send = st.form_submit_button("➤ Send")

b2, b3 = st.columns([1, 1])

with b2:
    ppt = st.button("📊 Create PPT", use_container_width=True)

with b3:
    clear = st.button("🗑 Clear Chat", use_container_width=True)

if send:
    if user_input.strip() == "":
        st.warning("Type something first.")
    elif st.session_state.mode is None:
        st.warning("Choose College or School mode first.")
    elif st.session_state.task is None:
        st.warning("Choose an option first.")
    else:
        with st.spinner("StudentMate AI thinking..."):
            prompt = f"Mode: {st.session_state.mode}. Task: {st.session_state.task}. User: {user_input}"

            reply = ask_ai(prompt)

            st.session_state.chat_history.append({"role": "user", "text": user_input})
            st.session_state.chat_history.append({"role": "bot", "text": reply})

            st.rerun()

if ppt:
    if st.session_state.mode is None:
        st.warning("Choose College or School mode first.")
    elif user_input.strip() == "":
        st.warning("Type PPT topic first.")
    else:
        with st.spinner("Creating PPT..."):
            ppt_file = create_ppt(user_input)

            with open(ppt_file, "rb") as f:
                st.download_button(
                    "Download PPT",
                    data=f,
                    file_name=ppt_file,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )

            st.success("PPT created successfully")

if clear:
    st.session_state.chat_history = []
    st.rerun()

st.markdown("</div>", unsafe_allow_html=True)
